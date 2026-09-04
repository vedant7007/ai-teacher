"""PDF/DOCX/PPTX -> Line records carrying page, font size, font name and text.

The hard part is PDFs like NCERT's, which fake bold by drawing the same heading
four or five times at sub-pixel offsets, each pass split into overlapping
horizontal fragments. Naive extraction yields "11.5 FA", "11.5 FACTORS ON WHICH
THE RESIST", "CTORS ON WHICH THE RESISTANCE OF A", "ANCE OF A" for one heading.

_cover_yband() fixes this by reconstructing each visual line as a greedy
left-to-right cover of its fragments: start at the leftmost, repeatedly take the
fragment that begins where the previous one ended and extends furthest right.
A fragment that starts well past the cursor begins a separate line, which keeps
side-boxes and multi-column layouts from being glued onto body text.
"""

from __future__ import annotations

from dataclasses import dataclass, asdict
from pathlib import Path

import pymupdf

# Fragments whose start is within this many points of the previous fragment's
# end are a continuation of the same visual line. Wider gaps start a new line.
JOIN_TOLERANCE_PT = 4.0
# Lines whose baselines are within this many points are the same visual row.
YBAND_TOLERANCE_PT = 1.5


@dataclass
class Line:
    text: str
    page: int          # 1-indexed
    size: float        # font size in points
    font: str
    y: float
    x0: float

    def to_dict(self) -> dict:
        return asdict(self)


def _cover_yband(frags: list[dict]) -> list[Line]:
    """Reconstruct visual lines from overlapping overprinted fragments."""
    frags = sorted(frags, key=lambda f: (f["x0"], -f["x1"]))
    out: list[Line] = []
    used: set[int] = set()

    while len(used) < len(frags):
        seed = next(i for i in range(len(frags)) if i not in used)
        # Widest fragment sharing the seed's left edge wins the start.
        start = max(
            (i for i, f in enumerate(frags)
             if i not in used and abs(f["x0"] - frags[seed]["x0"]) < JOIN_TOLERANCE_PT),
            key=lambda i: frags[i]["x1"],
        )
        parts = [frags[start]["text"]]
        cursor = frags[start]["x1"]
        used.add(start)
        # Anything sharing that left edge is an overprint duplicate, drop it.
        for i, f in enumerate(frags):
            if i not in used and abs(f["x0"] - frags[start]["x0"]) < JOIN_TOLERANCE_PT:
                used.add(i)

        while True:
            nxt = [
                i for i, f in enumerate(frags)
                if i not in used and abs(f["x0"] - cursor) < JOIN_TOLERANCE_PT
            ]
            if not nxt:
                break
            best = max(nxt, key=lambda i: frags[i]["x1"])
            parts.append(frags[best]["text"])
            cursor = frags[best]["x1"]
            for i in nxt:
                used.add(i)

        text = "".join(parts).strip()
        if text:
            f0 = frags[start]
            out.append(Line(text, f0["page"], f0["size"], f0["font"], f0["y"], f0["x0"]))

    # An overprint pass can start mid-word ("CTORS ON WHICH..."), so it never
    # joins the cover and survives as its own line. Any such leftover is a
    # substring of the line that was reconstructed properly.
    keep = [
        l for l in out
        if not any(l.text != o.text and l.text in o.text for o in out)
    ]
    return keep or out


def parse_pdf(path: str | Path) -> list[Line]:
    doc = pymupdf.open(str(path))
    lines: list[Line] = []
    for pno in range(len(doc)):
        frags: list[dict] = []
        for block in doc[pno].get_text("dict")["blocks"]:
            for line in block.get("lines", []):
                text = "".join(s["text"] for s in line["spans"])
                if not text.strip():
                    continue
                span = max(line["spans"], key=lambda s: len(s["text"]))
                x0, y0, x1, _ = line["bbox"]
                frags.append({
                    "text": text, "page": pno + 1, "size": round(span["size"], 1),
                    "font": span["font"], "y": y0, "x0": x0, "x1": x1,
                })
        # Group fragments into y-bands, then cover each band.
        for band in _yband_groups(frags):
            lines.extend(_cover_yband(band))
    doc.close()
    lines.sort(key=lambda l: (l.page, round(l.y, 0), l.x0))
    return lines


def _yband_groups(frags: list[dict]) -> list[list[dict]]:
    if not frags:
        return []
    groups: list[list[dict]] = []
    for f in sorted(frags, key=lambda f: f["y"]):
        if groups and abs(f["y"] - groups[-1][0]["y"]) <= YBAND_TOLERANCE_PT:
            groups[-1].append(f)
        else:
            groups.append([f])
    return groups


def parse_docx(path: str | Path) -> list[Line]:
    import docx

    d = docx.Document(str(path))
    lines = []
    for p in d.paragraphs:
        if not p.text.strip():
            continue
        # Word gives us the style directly, no clustering needed.
        size = 20.0 if p.style.name.startswith("Heading") else 10.5
        lines.append(Line(p.text.strip(), 1, size, p.style.name, 0.0, 0.0))
    return lines


def parse_pptx(path: str | Path) -> list[Line]:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if text:
                    size = 20.0 if para.level == 0 and len(text) < 80 else 10.5
                    lines.append(Line(text, i, size, "pptx", 0.0, 0.0))
    return lines


def parse(path: str | Path) -> list[Line]:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(path)
    if ext == ".docx":
        return parse_docx(path)
    if ext == ".pptx":
        return parse_pptx(path)
    if ext in {".txt", ".md"}:
        text = Path(path).read_text(encoding="utf-8", errors="replace")
        return [Line(l.strip(), 1, 10.5, "text", float(i), 0.0)
                for i, l in enumerate(text.splitlines()) if l.strip()]
    raise ValueError(f"unsupported file type: {ext}")
